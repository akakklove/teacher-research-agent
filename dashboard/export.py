"""
导出模块 — 科研报告导出为 PPT / Markdown
"""
from pathlib import Path
from datetime import datetime
from typing import Optional


def _get_val(m, key, default=None):
    """兼容 dict 和对象两种 metric 格式"""
    if isinstance(m, dict):
        return m.get(key, default)
    return getattr(m, key, default)


class ReportExporter:
    """报告导出器"""

    EXPORT_DIR = Path(__file__).parent.parent / "exports"

    # ── PPT 导出 ──
    @classmethod
    def to_pptx(cls, teacher, metrics, insights, narrative="",
                scenario_name="科研报告", time_range=("", ""), output_path=None):
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # 封面
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = RGBColor(0x0A, 0x0E, 0x27)

        cls._add_textbox(slide, Inches(1), Inches(2), Inches(11), Inches(1.5),
                         teacher.get('xm', '') + "老师", 44,
                         RGBColor(0xFF, 0xFF, 0xFF), PP_ALIGN.CENTER)
        cls._add_textbox(slide, Inches(1), Inches(3.5), Inches(11), Inches(1),
                         scenario_name, 28, RGBColor(0x09, 0x84, 0xE3), PP_ALIGN.CENTER)
        cls._add_textbox(slide, Inches(1), Inches(4.8), Inches(11), Inches(0.6),
                         f"{teacher.get('dw_mc','')} · {teacher.get('zc','')}",
                         18, RGBColor(0x63, 0x6E, 0x72), PP_ALIGN.CENTER)
        cls._add_textbox(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
                         f"数据范围: {time_range[0]} ~ {time_range[1]}",
                         14, RGBColor(0xB2, 0xBE, 0xC3), PP_ALIGN.CENTER)

        # KPI 概览页
        kpi_metrics = [m for m in metrics
                       if _get_val(m, "chart_type") == "kpi_card"
                       and _get_val(m, "value") is not None]
        if kpi_metrics:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            cls._add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
                             "核心指标概览", 28, RGBColor(0x09, 0x84, 0xE3))
            cols = 5
            for i, m in enumerate(kpi_metrics):
                row, col = i // cols, i % cols
                cls._add_kpi_card(slide,
                                  Inches(0.5 + col * 2.5), Inches(1.5 + row * 2.0),
                                  Inches(2.2), Inches(1.5),
                                  _get_val(m, "name", ""),
                                  _get_val(m, "value", 0),
                                  _get_val(m, "unit", ""))

        # 洞察页
        if insights or narrative:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            cls._add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
                             "AI 洞察", 28, RGBColor(0x09, 0x84, 0xE3))
            y = Inches(1.2)
            if insights:
                cls._add_textbox(slide, Inches(0.5), y, Inches(12), Inches(3),
                                 "\n\n".join(insights), 16, RGBColor(0x2D, 0x34, 0x36))
                y = Inches(4.5)
            if narrative:
                cls._add_textbox(slide, Inches(0.5), y, Inches(12), Inches(2.5),
                                 narrative[:500], 13, RGBColor(0x63, 0x6E, 0x72))

        # 保存
        cls.EXPORT_DIR.mkdir(parents=True, exist_ok=True)
        if output_path is None:
            ts = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_path = cls.EXPORT_DIR / f"{teacher.get('xm','report')}_{ts}.pptx"
        else:
            output_path = Path(output_path)
        prs.save(str(output_path))
        return str(output_path)

    # ── Markdown 导出 ──
    @classmethod
    def to_markdown(cls, teacher, metrics, insights, narrative="",
                    scenario_name="科研报告", output_path=None):
        name = teacher.get("xm", "老师")
        dept = teacher.get("dw_mc", "")
        title = teacher.get("zc", "")

        md = f"# {name}老师 — {scenario_name}\n\n"
        md += f"**单位：**{dept}　**职称：**{title}\n\n"
        md += f"**生成时间：**{datetime.now().strftime('%Y-%m-%d %H:%M')}\n\n---\n\n"

        # KPI
        kpis = [m for m in metrics
                if _get_val(m, "chart_type") == "kpi_card"
                and _get_val(m, "value") is not None]
        if kpis:
            md += "## 核心指标\n\n| 指标 | 数值 |\n|------|------|\n"
            for m in kpis:
                val = ReportExporter._format_kpi(_get_val(m, "value", 0), _get_val(m, "unit", ""))
                md += f"| {_get_val(m, 'name', '')} | {val} |\n"
            md += "\n"

        # 洞察
        if insights:
            md += "## AI 洞察\n\n"
            for ins in insights:
                md += f"- {ins}\n"
            md += "\n"

        # 叙事
        if narrative:
            md += "## 详细报告\n\n" + narrative + "\n"

        cls.EXPORT_DIR.mkdir(parents=True, exist_ok=True)
        if output_path is None:
            ts = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_path = cls.EXPORT_DIR / f"{teacher.get('xm','report')}_{ts}.md"
        else:
            output_path = Path(output_path)
        output_path.write_text(md, encoding="utf-8")
        return str(output_path)

    # ── 辅助方法 ──
    @staticmethod
    def _add_textbox(slide, left, top, width, height, text, size, color, align=None):
        from pptx.util import Pt
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        if align:
            p.alignment = align
        return txBox

    @staticmethod
    def _add_kpi_card(slide, left, top, width, height, label, value, unit):
        from pptx.util import Pt, Inches
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        shape = slide.shapes.add_shape(1, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x0D, 0x13, 0x37)
        shape.line.fill.background()

        val_text = ReportExporter._format_kpi(value, unit)
        txBox = slide.shapes.add_textbox(
            left + Inches(0.1), top + Inches(0.2), width - Inches(0.2), Inches(0.7))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = val_text
        p.font.size = Pt(28)
        p.font.color.rgb = RGBColor(0x09, 0x84, 0xE3)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        txBox2 = slide.shapes.add_textbox(
            left + Inches(0.1), top + Inches(0.9), width - Inches(0.2), Inches(0.4))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = label
        p2.font.size = Pt(11)
        p2.font.color.rgb = RGBColor(0xB2, 0xBE, 0xC3)
        p2.alignment = PP_ALIGN.CENTER

    @staticmethod
    def _format_kpi(value, unit):
        if value is None:
            return "-"
        if unit == "元" and abs(value) >= 10000:
            return f"{value/10000:.1f}万"
        if isinstance(value, float) and value == int(value):
            return str(int(value))
        if isinstance(value, float):
            return f"{value:.1f}"
        return str(value)
